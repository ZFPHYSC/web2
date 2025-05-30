import asyncio
from typing import List, Dict, Tuple, Optional
import os
import re
from pathlib import Path
import logging

# Document processing imports
import pdfplumber
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract

# Database imports
from sqlalchemy import select, update
from models.database import (
    AsyncSessionLocal, Document, DocumentChunk, 
    ProcessingQueue, Course
)

logger = logging.getLogger(__name__)

class IngestionService:
    def __init__(self):
        self.embedding_service = None
        self.text_extractors = {
            '.pdf': self.extract_pdf,
            '.docx': self.extract_docx,
            '.pptx': self.extract_pptx,
            '.xlsx': self.extract_xlsx,
            '.txt': self.extract_txt,
            '.jpg': self.extract_image,
            '.jpeg': self.extract_image,
            '.png': self.extract_image,
        }
        
    async def initialize(self, embedding_service):
        """Initialize with embedding service"""
        self.embedding_service = embedding_service
        logger.info("Ingestion service initialized")
    
    async def process_file(self, course_id: str, file_path: str, filename: str) -> bool:
        """Simplified file processing"""
        try:
            async with AsyncSessionLocal() as session:
                # Create document record
                document = Document(
                    course_id=course_id,
                    filename=filename,
                    original_path=file_path,
                    file_type=Path(filename).suffix.lower(),
                    file_size=os.path.getsize(file_path) if os.path.exists(file_path) else 0,
                    status="processing"
                )
                
                session.add(document)
                await session.flush()
                
                # Extract text
                text = await self.extract_text(file_path, filename)
                
                if not text or len(text.strip()) < 10:
                    document.status = "failed"
                    document.error_message = "No text extracted"
                    await session.commit()
                    return False
                
                document.raw_text = text
                
                # Simple chunking
                chunks = self.simple_chunk(text, chunk_size=1000, overlap=200)
                
                # Store chunks
                for i, chunk_text in enumerate(chunks):
                    chunk_record = DocumentChunk(
                        document_id=document.id,
                        course_id=course_id,
                        content=chunk_text,
                        chunk_index=i,
                        chunk_metadata={"filename": filename},
                        chunk_type="semantic"
                    )
                    session.add(chunk_record)
                
                await session.flush()
                
                # Store embeddings
                chunk_dicts = [{"content": c, "metadata": {"filename": filename}} for c in chunks]
                await self.embedding_service.store_embeddings(
                    chunk_dicts, 
                    course_id, 
                    str(document.id)
                )
                
                # Update document status
                document.status = "completed"
                document.chunk_count = len(chunks)
                
                # Update course file count
                await session.execute(
                    update(Course)
                    .where(Course.id == course_id)
                    .values(file_count=Course.file_count + 1)
                )
                
                await session.commit()
                
                # Clean up temp file
                if os.path.exists(file_path):
                    os.remove(file_path)
                
                logger.info(f"Processed {filename} with {len(chunks)} chunks")
                return True
                
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            return False
    
    async def extract_text(self, file_path: str, filename: str) -> str:
        """Extract text from a file based on its type"""
        try:
            file_ext = Path(filename).suffix.lower()
            
            if file_ext in self.text_extractors:
                extractor = self.text_extractors[file_ext]
                return await extractor(file_path)
            else:
                logger.warning(f"No extractor for file type: {file_ext}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            return ""
    
    async def extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF with fallback to OCR"""
        text = ""
        
        try:
            # Try pdfplumber first
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except Exception as e:
            logger.warning(f"PDFPlumber failed: {e}")
        
        # If no text extracted or very little, try OCR
        if len(text.strip()) < 100:
            try:
                text = await self.ocr_pdf(file_path)
            except Exception as e:
                logger.warning(f"OCR failed: {e}")
        
        return text
    
    async def ocr_pdf(self, file_path: str) -> str:
        """Extract text from PDF using OCR"""
        try:
            # Convert PDF to images and OCR
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            text = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                
                # Save as temporary image
                temp_img_path = f"temp_page_{page_num}.png"
                with open(temp_img_path, "wb") as f:
                    f.write(img_data)
                
                # OCR the image
                page_text = pytesseract.image_to_string(Image.open(temp_img_path))
                text += page_text + "\n\n"
                
                # Clean up
                os.remove(temp_img_path)
            
            doc.close()
            return text
            
        except Exception as e:
            logger.error(f"OCR error: {e}")
            return ""
    
    async def extract_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_path)
            text = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        text.append('\t'.join(row_text))
            
            return '\n'.join(text)
            
        except Exception as e:
            logger.error(f"Error extracting from DOCX: {e}")
            return ""
    
    async def extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides):
                text.append(f"\n--- Slide {slide_num + 1} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            
            return '\n'.join(text)
            
        except Exception as e:
            logger.error(f"Error extracting from PPTX: {e}")
            return ""
    
    async def extract_xlsx(self, file_path: str) -> str:
        """Extract text from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                text.append(f"\n--- Sheet: {sheet_name} ---\n")
                
                # Convert to string, handling NaN values
                sheet_text = df.to_string(index=False, na_rep='')
                text.append(sheet_text)
            
            return '\n'.join(text)
            
        except Exception as e:
            logger.error(f"Error extracting from XLSX: {e}")
            return ""
    
    async def extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error reading text file: {e}")
                return ""
    
    async def extract_image(self, file_path: str) -> str:
        """Extract text from images using OCR"""
        try:
            # Use PIL to open and process the image
            image = Image.open(file_path)
            
            # Convert to RGB if necessary (for better OCR results)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Use pytesseract to extract text
            text = pytesseract.image_to_string(image, config='--psm 6')
            
            # Clean up the text
            text = text.strip()
            
            if len(text) < 10:
                # If very little text found, try different PSM modes
                for psm in [3, 4, 6, 8, 11, 13]:
                    try:
                        alt_text = pytesseract.image_to_string(image, config=f'--psm {psm}')
                        if len(alt_text.strip()) > len(text):
                            text = alt_text.strip()
                    except:
                        continue
            
            return text if text else f"[Image file: {os.path.basename(file_path)}]"
            
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {e}")
            return f"[Image file: {os.path.basename(file_path)} - OCR failed]"
    
    def simple_chunk(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Simple text chunking"""
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < text_length:
                last_period = chunk.rfind('.')
                if last_period > chunk_size * 0.8:
                    end = start + last_period + 1
                    chunk = text[start:end]
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return [c for c in chunks if c]